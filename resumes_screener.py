#!/usr/bin/env python3
"""Resume Analyzer — On-Prem AI HR Screening Tool (Ollama-powered)."""

import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import threading
import os
import json
import re
import urllib.request
import urllib.error
from pathlib import Path

# ─── Text extraction ──────────────────────────────────────────────────────────

def extract_pdf(path: str) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    except Exception as exc:
        return f"[PDF read error: {exc}]"


def extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as exc:
        return f"[DOCX read error: {exc}]"


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as exc:
        return f"[PPTX read error: {exc}]"


def extract_txt(path: str) -> str:
    try:
        return Path(path).read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return f"[TXT read error: {exc}]"


def extract_text(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext in {".doc", ".docx"}:
        return extract_docx(path)
    if ext in {".ppt", ".pptx"}:
        return extract_pptx(path)
    if ext == ".txt":
        return extract_txt(path)
    return f"[Unsupported file type: {ext}]"


# ─── Ollama interaction ───────────────────────────────────────────────────────

ANALYSIS_PROMPT = """\
You are a seasoned HR analyst with more than 20 years of recruitment experience
across the industry. You have screened thousands of resumes, interviewed candidates
at every level, and have a sharp eye for role fit, career trajectory, and red flags.
Apply that depth of judgement when evaluating this resume against the job description.

Respond with ONLY a valid JSON object — no markdown fences, no explanation, just raw JSON.

JOB DESCRIPTION:
{jd}

RESUME:
{resume}

Return this exact JSON (fill in real values):
{{
  "name": "<candidate full name>",
  "education": [
    {{"degree": "<e.g., B.Tech>", "college": "<e.g., NIT Warangal>", "year": <year of passing as int, e.g., 2018>}}
  ],
  "current_organization": "<current employer, or empty string if not currently employed>",
  "past_organizations": ["<previous employer 1>", "<previous employer 2>"],
  "total_experience": {{"years": <int>, "months": <int 0-11>}},
  "relevant_experience": {{"years": <int>, "months": <int 0-11>}},
  "highlights": ["<strength 1>", "<strength 2>", "<strength 3>"],
  "lowlights": ["<gap 1>", "<gap 2>"],
  "score": <integer 1-10>,
  "decision": "<Select if score > 6, else Reject>"
}}

Rules:
- education: include ONLY degree, college name, and year of passing.
  Do NOT include location, GPA, percentage, specialisation, coursework, or any other detail.
  List every distinct qualification (e.g., both B.Tech and M.Tech).
  Do NOT repeat the same qualification more than once.
- current_organization: the single employer the candidate works at right now.
  Use an empty string if the candidate is not currently employed.
- past_organizations: every prior employer in reverse chronological order (most recent first).
  Deduplicate strictly — each company name must appear only once even if the candidate
  worked there in multiple stints. Do NOT include the current organization in this list.
- score is 1–10 based on how well the resume matches the job description
- decision MUST be "Select" when score > 6, otherwise "Reject"
- Extract the candidate name directly from the resume text
- Be specific in highlights and lowlights (3 highlights, 2 lowlights minimum)
"""

OLLAMA_TIMEOUT = 600  # seconds — local LLMs can be slow on CPU


def _strip_fences(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^```[a-z]*\s*", "", text)
    text = re.sub(r"\s*```$", "", text)
    return text.strip()


def _ollama_url(host: str, port: str, path: str) -> str:
    host = host.strip().rstrip("/")
    if not host.startswith(("http://", "https://")):
        host = f"http://{host}"
    return f"{host}:{port}{path}"


def ollama_generate(host: str, port: str, model: str, prompt: str) -> str:
    """Call Ollama /api/generate and return the model's text response."""
    url = _ollama_url(host, port, "/api/generate")
    payload = json.dumps({
        "model": model,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": {"temperature": 0.2},
    }).encode("utf-8")
    req = urllib.request.Request(
        url, data=payload, headers={"Content-Type": "application/json"}, method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=OLLAMA_TIMEOUT) as resp:
            data = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"Ollama HTTP {e.code}: {body}") from e
    except urllib.error.URLError as e:
        raise RuntimeError(f"Cannot reach Ollama at {url} — {e.reason}") from e
    return data.get("response", "")


def ollama_list_models(host: str, port: str) -> list[str]:
    """Return locally available model names from Ollama, or [] on failure."""
    url = _ollama_url(host, port, "/api/tags")
    try:
        with urllib.request.urlopen(url, timeout=10) as resp:
            data = json.loads(resp.read().decode("utf-8"))
        return [m.get("name", "") for m in data.get("models", []) if m.get("name")]
    except Exception:
        return []


def analyse_resume(path: str, jd_text: str, host: str, port: str, model: str) -> dict:
    resume_text = extract_text(path)
    prompt = ANALYSIS_PROMPT.format(jd=jd_text, resume=resume_text)
    raw = ollama_generate(host, port, model, prompt)
    return json.loads(_strip_fences(raw))


# ─── Design tokens ────────────────────────────────────────────────────────────

# Slate + Indigo professional palette
PRIMARY       = "#0F172A"   # slate-900   (header background)
PRIMARY_SOFT  = "#1E293B"   # slate-800
ACCENT        = "#4F46E5"   # indigo-600  (primary action)
ACCENT_HOVER  = "#4338CA"   # indigo-700
ACCENT_DEEP   = "#3730A3"   # indigo-800
ACCENT_LIGHT  = "#EEF2FF"   # indigo-50
BG            = "#F1F5F9"   # slate-100
CARD          = "#FFFFFF"
BORDER        = "#E2E8F0"   # slate-200
BORDER_DARK   = "#CBD5E1"   # slate-300
TEXT          = "#0F172A"
TEXT_MUTED    = "#475569"   # slate-600
TEXT_SUBTLE   = "#94A3B8"   # slate-400
SUCCESS       = "#047857"   # emerald-700
SUCCESS_BG    = "#D1FAE5"   # emerald-100
DANGER        = "#B91C1C"   # red-700
DANGER_BG     = "#FEE2E2"   # red-100
NEUTRAL_BTN   = "#F1F5F9"
NEUTRAL_HOVER = "#E2E8F0"
ROW_ALT       = "#F8FAFC"   # slate-50

FONT_FAMILY  = "Segoe UI"
FONT_TITLE   = (FONT_FAMILY, 17, "bold")
FONT_SUB     = (FONT_FAMILY, 10)
FONT_SECTION = (FONT_FAMILY, 11, "bold")
FONT_BODY    = (FONT_FAMILY, 10)
FONT_BOLD    = (FONT_FAMILY, 10, "bold")
FONT_LABEL   = (FONT_FAMILY, 9, "bold")
FONT_SMALL   = (FONT_FAMILY, 9)
FONT_MICRO   = (FONT_FAMILY, 8)
FONT_MONO    = ("Consolas", 10)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def fmt_exp(exp: dict | None) -> str:
    if not exp:
        return "N/A"
    y, m = int(exp.get("years", 0)), int(exp.get("months", 0))
    return f"{y}y {m}m"


def fmt_education(edu: list | None) -> str:
    """Render education as 'Degree, College, Year' entries, deduplicated."""
    if not edu:
        return "—"
    seen, out = set(), []
    for e in edu:
        if not isinstance(e, dict):
            continue
        deg  = str(e.get("degree", "") or "").strip()
        coll = str(e.get("college", "") or "").strip()
        yr   = e.get("year")
        yr_s = str(int(yr)) if isinstance(yr, (int, float)) and yr else (
            str(yr).strip() if yr else ""
        )
        parts = [p for p in (deg, coll, yr_s) if p]
        if not parts:
            continue
        line = ", ".join(parts)
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(line)
    return "  ·  ".join(out) if out else "—"


def fmt_orgs(orgs: list | None, exclude: str = "") -> str:
    """Comma-separated list of organizations, deduplicated, excluding `exclude`."""
    if not orgs:
        return "—"
    excl = (exclude or "").strip().lower()
    seen, out = set(), []
    for o in orgs:
        s = str(o or "").strip()
        if not s:
            continue
        key = s.lower()
        if key in seen or key == excl:
            continue
        seen.add(key)
        out.append(s)
    return ", ".join(out) if out else "—"


class ToolTip:
    def __init__(self, widget, text: str):
        self.tip = None
        widget.bind("<Enter>", lambda _: self._show(widget, text))
        widget.bind("<Leave>", lambda _: self._hide())

    def _show(self, widget, text):
        x = widget.winfo_rootx() + 20
        y = widget.winfo_rooty() + widget.winfo_height() + 6
        self.tip = tk.Toplevel(widget)
        self.tip.wm_overrideredirect(True)
        self.tip.wm_geometry(f"+{x}+{y}")
        tk.Label(
            self.tip, text=text, bg=PRIMARY_SOFT, fg="white",
            relief="flat", borderwidth=0, font=FONT_SMALL,
            padx=8, pady=5,
        ).pack()

    def _hide(self):
        if self.tip:
            self.tip.destroy()
            self.tip = None


class HoverButton(tk.Button):
    """Flat button with controlled hover/press colors."""
    def __init__(self, parent, text, command, *,
                 bg=ACCENT, fg="white", hover=ACCENT_HOVER,
                 font=FONT_BODY, padx=14, pady=7, **kw):
        super().__init__(
            parent, text=text, command=command, font=font,
            bg=bg, fg=fg, activebackground=hover, activeforeground=fg,
            relief="flat", borderwidth=0, padx=padx, pady=pady,
            cursor="hand2", highlightthickness=0, **kw,
        )
        self._bg, self._hv = bg, hover
        self.bind("<Enter>", lambda _: self.config(bg=self._hv))
        self.bind("<Leave>", lambda _: self.config(bg=self._bg))


# ─── Main application ─────────────────────────────────────────────────────────

class ResumeAnalyzerApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Resume Analyzer — On-Prem AI HR Tool")
        self.geometry("1320x860")
        self.minsize(1000, 660)
        self.configure(bg=BG)
        self.resume_files: list[str] = []
        self.results: list[dict] = []
        self._apply_style()
        self._build_ui()

    # ── Style ─────────────────────────────────────────────────────────────────

    def _apply_style(self):
        s = ttk.Style(self)
        s.theme_use("clam")

        s.configure(
            "TEntry",
            fieldbackground="white", background="white",
            bordercolor=BORDER, lightcolor=BORDER, darkcolor=BORDER,
            relief="flat", padding=8, foreground=TEXT,
        )
        s.map(
            "TEntry",
            bordercolor=[("focus", ACCENT)],
            lightcolor=[("focus", ACCENT)],
            darkcolor=[("focus", ACCENT)],
        )

        s.configure(
            "TCombobox",
            fieldbackground="white", background="white",
            bordercolor=BORDER, lightcolor=BORDER, darkcolor=BORDER,
            relief="flat", padding=6, foreground=TEXT,
            arrowcolor=TEXT_MUTED,
        )
        s.map("TCombobox",
              fieldbackground=[("readonly", "white")],
              bordercolor=[("focus", ACCENT)])

        s.configure("Vertical.TScrollbar",
                    troughcolor=BG, background=BORDER_DARK,
                    bordercolor=BG, arrowcolor=TEXT_MUTED,
                    relief="flat")
        s.configure("Horizontal.TScrollbar",
                    troughcolor=BG, background=BORDER_DARK,
                    bordercolor=BG, arrowcolor=TEXT_MUTED,
                    relief="flat")

        s.configure(
            "Results.Treeview",
            rowheight=58, font=FONT_SMALL,
            background="white", fieldbackground="white",
            foreground=TEXT, bordercolor=BORDER,
            relief="flat",
        )
        s.configure(
            "Results.Treeview.Heading",
            font=FONT_BOLD, background=PRIMARY, foreground="white",
            relief="flat", padding=10, borderwidth=0,
        )
        s.map("Results.Treeview.Heading",
              background=[("active", PRIMARY_SOFT)])
        s.map("Results.Treeview",
              background=[("selected", ACCENT_LIGHT)],
              foreground=[("selected", TEXT)])

        s.configure(
            "Modern.Horizontal.TProgressbar",
            troughcolor=BORDER, background=ACCENT,
            bordercolor=BORDER, lightcolor=ACCENT, darkcolor=ACCENT,
            thickness=8,
        )

    # ── Layout ────────────────────────────────────────────────────────────────

    def _build_ui(self):
        self._build_header()
        self._build_scroll_area()

        self._build_config_card()
        self._build_jd_card()
        self._build_resumes_card()
        self._build_action_row()
        self._build_results_card()

        self._build_footer()

    def _build_header(self):
        hdr = tk.Frame(self, bg=PRIMARY, height=72)
        hdr.pack(fill="x", side="top")
        hdr.pack_propagate(False)

        left = tk.Frame(hdr, bg=PRIMARY)
        left.pack(side="left", padx=24, fill="y")

        tk.Label(left, text="◆", font=(FONT_FAMILY, 22, "bold"),
                 bg=PRIMARY, fg=ACCENT).pack(side="left", pady=18, padx=(0, 10))

        title_col = tk.Frame(left, bg=PRIMARY)
        title_col.pack(side="left", pady=14)
        tk.Label(title_col, text="Resume Analyzer", font=FONT_TITLE,
                 bg=PRIMARY, fg="white").pack(anchor="w")
        tk.Label(title_col, text="On-Prem AI HR Screening · Ollama",
                 font=FONT_SUB, bg=PRIMARY, fg=TEXT_SUBTLE).pack(anchor="w")

        right = tk.Frame(hdr, bg=PRIMARY)
        right.pack(side="right", padx=24, fill="y")

        self._status_dot = tk.Label(right, text="●", font=(FONT_FAMILY, 12),
                                    bg=PRIMARY, fg=TEXT_SUBTLE)
        self._status_dot.pack(side="left", pady=24)
        self._status_lbl = tk.Label(right, text="Not connected",
                                    font=FONT_SMALL, bg=PRIMARY, fg=TEXT_SUBTLE)
        self._status_lbl.pack(side="left", padx=(6, 0), pady=24)

    def _build_scroll_area(self):
        outer = tk.Frame(self, bg=BG)
        outer.pack(fill="both", expand=True)

        canvas = tk.Canvas(outer, bg=BG, highlightthickness=0)
        vscroll = ttk.Scrollbar(outer, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=vscroll.set)
        vscroll.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        self._main = tk.Frame(canvas, bg=BG)
        win_id = canvas.create_window((0, 0), window=self._main, anchor="nw")

        def _on_configure(_e):
            canvas.configure(scrollregion=canvas.bbox("all"))
        def _on_canvas_resize(e):
            canvas.itemconfig(win_id, width=e.width)

        self._main.bind("<Configure>", _on_configure)
        canvas.bind("<Configure>", _on_canvas_resize)
        canvas.bind_all("<MouseWheel>",
                        lambda e: canvas.yview_scroll(-1*(e.delta//120), "units"))

    def _build_footer(self):
        ft = tk.Frame(self, bg=PRIMARY_SOFT, height=28)
        ft.pack(fill="x", side="bottom")
        ft.pack_propagate(False)
        tk.Label(ft, text="On-prem · No data leaves your network",
                 font=FONT_MICRO, bg=PRIMARY_SOFT,
                 fg=TEXT_SUBTLE).pack(side="left", padx=24, pady=6)
        tk.Label(ft, text="v2.0 · Ollama backend",
                 font=FONT_MICRO, bg=PRIMARY_SOFT,
                 fg=TEXT_SUBTLE).pack(side="right", padx=24, pady=6)

    # ── Card primitives ───────────────────────────────────────────────────────

    def _card(self, title: str, subtitle: str = "") -> tk.Frame:
        # Outer creates a soft "shadow" border feel via two-tone padding
        wrapper = tk.Frame(self._main, bg=BG)
        wrapper.pack(fill="x", padx=24, pady=10)

        outer = tk.Frame(wrapper, bg=CARD,
                         highlightbackground=BORDER, highlightthickness=1)
        outer.pack(fill="x")

        head = tk.Frame(outer, bg=CARD)
        head.pack(fill="x", padx=20, pady=(14, 4))
        tk.Label(head, text=title, font=FONT_SECTION,
                 bg=CARD, fg=TEXT).pack(anchor="w")
        if subtitle:
            tk.Label(head, text=subtitle, font=FONT_SMALL,
                     bg=CARD, fg=TEXT_MUTED).pack(anchor="w", pady=(2, 0))

        sep = tk.Frame(outer, bg=BORDER, height=1)
        sep.pack(fill="x", padx=20, pady=(8, 0))

        body = tk.Frame(outer, bg=CARD)
        body.pack(fill="x", padx=20, pady=(10, 16))
        return body

    def _label(self, parent, text):
        return tk.Label(parent, text=text, font=FONT_LABEL,
                        bg=CARD, fg=TEXT_MUTED)

    # ── Config card (Ollama) ──────────────────────────────────────────────────

    def _build_config_card(self):
        body = self._card(
            "Ollama Connection",
            "Configure the on-prem Ollama endpoint serving your enterprise LLM.",
        )

        grid = tk.Frame(body, bg=CARD)
        grid.pack(fill="x")

        # Row 1: host, port, test
        self._label(grid, "HOST / IP").grid(row=0, column=0, sticky="w", padx=(0, 6))
        self._label(grid, "PORT").grid(row=0, column=1, sticky="w", padx=6)
        self._label(grid, "MODEL").grid(row=0, column=2, sticky="w", padx=6)

        self.host_var  = tk.StringVar(value="localhost")
        self.port_var  = tk.StringVar(value="11434")
        self.model_var = tk.StringVar(value="llama3.1")

        host_e = ttk.Entry(grid, textvariable=self.host_var, width=22, font=FONT_BODY)
        host_e.grid(row=1, column=0, sticky="we", padx=(0, 6), pady=(2, 0))
        ToolTip(host_e, "Examples:  localhost  ·  10.0.0.42  ·  ollama.intra")

        port_e = ttk.Entry(grid, textvariable=self.port_var, width=8, font=FONT_BODY)
        port_e.grid(row=1, column=1, sticky="we", padx=6, pady=(2, 0))
        ToolTip(port_e, "Default Ollama port is 11434")

        self._model_combo = ttk.Combobox(
            grid, textvariable=self.model_var, font=FONT_BODY,
            values=["llama3.1", "llama3", "mistral", "qwen2.5", "phi3"],
            width=24,
        )
        self._model_combo.grid(row=1, column=2, sticky="we", padx=6, pady=(2, 0))
        ToolTip(self._model_combo, "Type the exact model name pulled in Ollama,\n"
                                   "or click Refresh to fetch the list.")

        btns = tk.Frame(grid, bg=CARD)
        btns.grid(row=1, column=3, sticky="w", padx=(10, 0), pady=(2, 0))

        HoverButton(btns, "↻  Refresh Models",
                    self._refresh_models, bg=NEUTRAL_BTN, fg=TEXT,
                    hover=NEUTRAL_HOVER, font=FONT_SMALL,
                    padx=10, pady=6).pack(side="left", padx=(0, 6))
        HoverButton(btns, "✓  Test Connection",
                    self._test_connection, bg=ACCENT, hover=ACCENT_HOVER,
                    font=FONT_BOLD, padx=14, pady=6).pack(side="left")

        for col, w in enumerate([3, 1, 3, 4]):
            grid.grid_columnconfigure(col, weight=w)

        # Hint line
        hint = tk.Frame(body, bg=CARD)
        hint.pack(fill="x", pady=(12, 0))
        tk.Label(hint, text="🔒",
                 font=(FONT_FAMILY, 10), bg=CARD,
                 fg=ACCENT).pack(side="left", padx=(0, 6))
        tk.Label(hint,
                 text="All inference happens inside your network. "
                      "No API keys, no external calls.",
                 font=FONT_SMALL, bg=CARD,
                 fg=TEXT_MUTED).pack(side="left")

    # ── JD card ───────────────────────────────────────────────────────────────

    def _build_jd_card(self):
        body = self._card(
            "Job Description",
            "Pick the role spec to evaluate every resume against.",
        )

        row = tk.Frame(body, bg=CARD)
        row.pack(fill="x")

        self.jd_var = tk.StringVar()
        jd_entry = ttk.Entry(row, textvariable=self.jd_var, font=FONT_BODY)
        jd_entry.pack(side="left", fill="x", expand=True, padx=(0, 10), ipady=2)
        ToolTip(jd_entry, "Supported formats: PDF · DOCX · PPTX · TXT")

        HoverButton(row, "📂  Browse File", self._browse_jd,
                    bg=NEUTRAL_BTN, fg=TEXT, hover=NEUTRAL_HOVER,
                    font=FONT_BOLD, padx=14, pady=7).pack(side="left")

    # ── Resumes card ──────────────────────────────────────────────────────────

    def _build_resumes_card(self):
        body = self._card(
            "Resumes",
            "Add the candidate documents you want to screen.",
        )

        btn_row = tk.Frame(body, bg=CARD)
        btn_row.pack(fill="x", pady=(0, 10))

        HoverButton(btn_row, "＋  Add Resumes", self._add_resumes,
                    bg=ACCENT, hover=ACCENT_HOVER,
                    font=FONT_BOLD, padx=14, pady=7).pack(side="left", padx=(0, 8))
        HoverButton(btn_row, "−  Remove Selected", self._remove_selected,
                    bg=NEUTRAL_BTN, fg=TEXT, hover=NEUTRAL_HOVER,
                    font=FONT_SMALL, padx=12, pady=7).pack(side="left", padx=(0, 8))
        HoverButton(btn_row, "✕  Clear All", self._clear_resumes,
                    bg=NEUTRAL_BTN, fg=TEXT, hover=NEUTRAL_HOVER,
                    font=FONT_SMALL, padx=12, pady=7).pack(side="left")

        self._count_pill = tk.Label(
            btn_row, text="0 files",
            font=FONT_BOLD, bg=ACCENT_LIGHT, fg=ACCENT_DEEP,
            padx=10, pady=4,
        )
        self._count_pill.pack(side="right")

        list_wrap = tk.Frame(body, bg=BORDER, padx=1, pady=1)
        list_wrap.pack(fill="x")
        list_inner = tk.Frame(list_wrap, bg=CARD)
        list_inner.pack(fill="both", expand=True)

        self._listbox = tk.Listbox(
            list_inner, height=6, font=FONT_BODY,
            bg=CARD, fg=TEXT,
            relief="flat", borderwidth=0,
            selectmode="extended", activestyle="none",
            selectbackground=ACCENT_LIGHT, selectforeground=TEXT,
            highlightthickness=0,
        )
        sb = ttk.Scrollbar(list_inner, orient="vertical",
                           command=self._listbox.yview)
        self._listbox.configure(yscrollcommand=sb.set)
        self._listbox.pack(side="left", fill="both", expand=True, padx=4, pady=4)
        sb.pack(side="right", fill="y")

    # ── Action row ────────────────────────────────────────────────────────────

    def _build_action_row(self):
        wrapper = tk.Frame(self._main, bg=BG)
        wrapper.pack(fill="x", padx=24, pady=(4, 10))

        card = tk.Frame(wrapper, bg=CARD,
                        highlightbackground=BORDER, highlightthickness=1)
        card.pack(fill="x")

        inner = tk.Frame(card, bg=CARD)
        inner.pack(fill="x", padx=20, pady=14)

        self._analyse_btn = HoverButton(
            inner, "▶  Run Analysis", self._start_analysis,
            bg=ACCENT, hover=ACCENT_HOVER,
            font=(FONT_FAMILY, 12, "bold"),
            padx=28, pady=12,
        )
        self._analyse_btn.pack(side="left", padx=(0, 18))

        progress_col = tk.Frame(inner, bg=CARD)
        progress_col.pack(side="left", fill="x", expand=True)

        self._progress_lbl = tk.Label(progress_col, text="Idle — ready to analyse",
                                      font=FONT_SMALL, bg=CARD, fg=TEXT_MUTED)
        self._progress_lbl.pack(anchor="w")

        self._progress = ttk.Progressbar(
            progress_col, mode="determinate", length=600,
            style="Modern.Horizontal.TProgressbar",
        )
        self._progress.pack(anchor="w", pady=(6, 0), fill="x", expand=True)

    # ── Results card ──────────────────────────────────────────────────────────

    def _build_results_card(self):
        wrapper = tk.Frame(self._main, bg=BG)
        wrapper.pack(fill="both", expand=True, padx=24, pady=(0, 18))

        card = tk.Frame(wrapper, bg=CARD,
                        highlightbackground=BORDER, highlightthickness=1)
        card.pack(fill="both", expand=True)

        head = tk.Frame(card, bg=CARD)
        head.pack(fill="x", padx=20, pady=(14, 4))
        tk.Label(head, text="Analysis Results", font=FONT_SECTION,
                 bg=CARD, fg=TEXT).pack(anchor="w")
        tk.Label(head, text="Double-click a row to view full candidate detail.",
                 font=FONT_SMALL, bg=CARD, fg=TEXT_MUTED).pack(anchor="w",
                                                               pady=(2, 0))

        sep = tk.Frame(card, bg=BORDER, height=1)
        sep.pack(fill="x", padx=20, pady=(8, 0))

        export_row = tk.Frame(card, bg=CARD)
        export_row.pack(fill="x", padx=20, pady=10)

        self._summary_lbl = tk.Label(
            export_row, text="No analysis run yet",
            font=FONT_SMALL, bg=CARD, fg=TEXT_MUTED,
        )
        self._summary_lbl.pack(side="left")

        HoverButton(export_row, "⬇  Export Excel",
                    self._export_excel,
                    bg=ACCENT, hover=ACCENT_HOVER,
                    font=FONT_BOLD, padx=14, pady=7
                    ).pack(side="right", padx=(8, 0))
        HoverButton(export_row, "⬇  Export CSV",
                    self._export_csv,
                    bg=NEUTRAL_BTN, fg=TEXT, hover=NEUTRAL_HOVER,
                    font=FONT_SMALL, padx=12, pady=7
                    ).pack(side="right")

        cols = ("#", "Name", "Education", "Current Org.", "Past Orgs.",
                "Total Exp.", "Relevant Exp.",
                "Highlights", "Lowlights", "Score", "Decision")
        widths = (40, 160, 240, 150, 220, 90, 100, 280, 220, 70, 100)

        tframe = tk.Frame(card, bg=CARD)
        tframe.pack(fill="both", expand=True, padx=20, pady=(0, 16))

        self.tree = ttk.Treeview(
            tframe, columns=cols, show="headings",
            style="Results.Treeview", selectmode="browse",
        )
        for col, w in zip(cols, widths):
            anchor = "center" if col in ("#", "Score", "Decision") else "w"
            self.tree.heading(col, text=col)
            self.tree.column(
                col, width=w, anchor=anchor,
                stretch=(col in ("Education", "Past Orgs.",
                                 "Highlights", "Lowlights")),
            )

        vsb = ttk.Scrollbar(tframe, orient="vertical", command=self.tree.yview)
        hsb = ttk.Scrollbar(tframe, orient="horizontal", command=self.tree.xview)
        self.tree.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)
        self.tree.grid(row=0, column=0, sticky="nsew")
        vsb.grid(row=0, column=1, sticky="ns")
        hsb.grid(row=1, column=0, sticky="ew")
        tframe.grid_rowconfigure(0, weight=1)
        tframe.grid_columnconfigure(0, weight=1)

        self.tree.tag_configure("select", background=SUCCESS_BG, foreground=TEXT)
        self.tree.tag_configure("reject", background=DANGER_BG,  foreground=TEXT)
        self.tree.tag_configure("even",   background=ROW_ALT)

        self.tree.bind("<Double-1>", self._show_detail)

    # ── Connection helpers ────────────────────────────────────────────────────

    def _set_status(self, ok: bool | None, text: str):
        color = TEXT_SUBTLE if ok is None else (
            "#22C55E" if ok else "#EF4444"
        )
        self._status_dot.config(fg=color)
        self._status_lbl.config(text=text, fg="white" if ok else TEXT_SUBTLE)

    def _refresh_models(self):
        host, port = self.host_var.get(), self.port_var.get()
        models = ollama_list_models(host, port)
        if models:
            self._model_combo["values"] = models
            if self.model_var.get() not in models:
                self.model_var.set(models[0])
            self._set_status(True, f"Connected · {len(models)} model(s)")
            messagebox.showinfo(
                "Models Refreshed",
                f"Found {len(models)} model(s) on {host}:{port}:\n\n"
                + "\n".join(f"• {m}" for m in models),
            )
        else:
            self._set_status(False, "Unreachable")
            messagebox.showwarning(
                "No Models",
                f"Could not reach Ollama at {host}:{port}, or no models pulled.\n\n"
                "Verify the host/port and that `ollama serve` is running.",
            )

    def _test_connection(self):
        host, port, model = (
            self.host_var.get(), self.port_var.get(), self.model_var.get()
        )
        models = ollama_list_models(host, port)
        if not models:
            self._set_status(False, "Unreachable")
            messagebox.showerror(
                "Connection Failed",
                f"Cannot reach Ollama at {host}:{port}.\n\n"
                "Check the address, firewall rules, and that the Ollama "
                "service is running on the target host.",
            )
            return
        if model not in models:
            self._set_status(True, f"Connected · model missing")
            messagebox.showwarning(
                "Model Not Found",
                f"Connected to Ollama at {host}:{port}, but model "
                f"'{model}' is not pulled.\n\nAvailable models:\n"
                + "\n".join(f"• {m}" for m in models)
                + f"\n\nRun on the server:\n  ollama pull {model}",
            )
            return
        self._set_status(True, f"Connected · {model}")
        messagebox.showinfo(
            "Connection OK",
            f"Successfully connected to Ollama at {host}:{port}.\n\n"
            f"Model '{model}' is ready.",
        )

    # ── Event handlers ────────────────────────────────────────────────────────

    def _browse_jd(self):
        path = filedialog.askopenfilename(
            title="Select Job Description File",
            filetypes=[("Supported Documents",
                        "*.pdf *.doc *.docx *.ppt *.pptx *.txt"),
                       ("All files", "*.*")],
        )
        if path:
            self.jd_var.set(path)

    def _add_resumes(self):
        paths = filedialog.askopenfilenames(
            title="Select Resume Files",
            filetypes=[("Supported Documents",
                        "*.pdf *.doc *.docx *.ppt *.pptx *.txt"),
                       ("All files", "*.*")],
        )
        for p in paths:
            if p not in self.resume_files:
                self.resume_files.append(p)
                self._listbox.insert("end", "  " + os.path.basename(p))
        self._update_count()

    def _clear_resumes(self):
        self.resume_files.clear()
        self._listbox.delete(0, "end")
        self._update_count()

    def _remove_selected(self):
        for idx in reversed(self._listbox.curselection()):
            self._listbox.delete(idx)
            self.resume_files.pop(idx)
        self._update_count()

    def _update_count(self):
        n = len(self.resume_files)
        self._count_pill.config(text=f"{n} file{'s' if n != 1 else ''}")

    def _validate(self) -> bool:
        if not self.host_var.get().strip():
            messagebox.showwarning("Missing Input", "Please enter the Ollama host or IP.")
            return False
        if not self.port_var.get().strip().isdigit():
            messagebox.showwarning("Invalid Input", "Port must be a number.")
            return False
        if not self.model_var.get().strip():
            messagebox.showwarning("Missing Input", "Please enter the Ollama model name.")
            return False
        if not self.jd_var.get():
            messagebox.showwarning("Missing Input", "Please select a Job Description file.")
            return False
        if not os.path.isfile(self.jd_var.get()):
            messagebox.showwarning("File Not Found",
                                   f"Job Description file not found:\n{self.jd_var.get()}")
            return False
        if not self.resume_files:
            messagebox.showwarning("Missing Input",
                                   "Please add at least one resume file.")
            return False
        return True

    def _start_analysis(self):
        if not self._validate():
            return
        self._analyse_btn.config(state="disabled")
        self.results.clear()
        for iid in self.tree.get_children():
            self.tree.delete(iid)
        self._summary_lbl.config(text="Analysis in progress…")
        threading.Thread(target=self._run_analysis, daemon=True).start()

    def _run_analysis(self):
        try:
            jd_text = extract_text(self.jd_var.get())
            host    = self.host_var.get().strip()
            port    = self.port_var.get().strip()
            model   = self.model_var.get().strip()
            total   = len(self.resume_files)

            for i, path in enumerate(self.resume_files):
                fname = os.path.basename(path)
                self._update_progress(i, total, f"Analysing  {fname}  …")
                try:
                    result = analyse_resume(path, jd_text, host, port, model)
                    result["_file"] = fname
                    self.results.append(result)
                    self._append_row(len(self.results), result)
                except json.JSONDecodeError as exc:
                    self._err(fname, f"Model returned invalid JSON.\n\nDetails: {exc}")
                except Exception as exc:
                    self._err(fname, str(exc))

            self._update_progress(total, total,
                                  f"Complete — {total} resume(s) analysed.")
            sel  = sum(1 for r in self.results
                       if str(r.get("decision", "")).strip().lower() == "select")
            rej  = len(self.results) - sel
            self.after(0, lambda: self._summary_lbl.config(
                text=f"{len(self.results)} analysed  ·  "
                     f"{sel} selected  ·  {rej} rejected"
            ))
        finally:
            self.after(0, lambda: self._analyse_btn.config(state="normal"))

    def _update_progress(self, current: int, total: int, msg: str):
        pct = (current / total * 100) if total else 0
        def _do():
            self._progress.config(value=pct)
            self._progress_lbl.config(text=msg)
        self.after(0, _do)

    def _append_row(self, serial: int, r: dict):
        highlights = "  •  ".join(r.get("highlights", []))
        lowlights  = "  •  ".join(r.get("lowlights",  []))
        score      = r.get("score", 0)
        decision   = r.get("decision", "Reject")
        tag        = "select" if str(decision).strip().lower() == "select" else "reject"

        current_org = str(r.get("current_organization") or "").strip() or "—"

        values = (
            serial,
            r.get("name", "Unknown"),
            fmt_education(r.get("education")),
            current_org,
            fmt_orgs(r.get("past_organizations"), exclude=current_org),
            fmt_exp(r.get("total_experience")),
            fmt_exp(r.get("relevant_experience")),
            highlights,
            lowlights,
            f"{score} / 10",
            "✓ " + decision if tag == "select" else "✕ " + decision,
        )
        self.after(0, lambda v=values, t=tag: self.tree.insert("", "end", values=v, tags=(t,)))

    def _err(self, fname: str, msg: str):
        self.after(0, lambda: messagebox.showerror(
            "Analysis Error", f"Failed to analyse:\n{fname}\n\n{msg}"
        ))

    # ── Detail popup ──────────────────────────────────────────────────────────

    def _show_detail(self, _event=None):
        sel = self.tree.selection()
        if not sel:
            return
        idx = self.tree.index(sel[0])
        if idx >= len(self.results):
            return
        r = self.results[idx]

        win = tk.Toplevel(self)
        win.title(f"Detail — {r.get('name', 'Unknown')}")
        win.geometry("760x780")
        win.configure(bg=BG)
        win.grab_set()
        win.resizable(True, True)

        # Header
        hdr = tk.Frame(win, bg=PRIMARY, height=66)
        hdr.pack(fill="x")
        hdr.pack_propagate(False)

        hdr_inner = tk.Frame(hdr, bg=PRIMARY)
        hdr_inner.pack(side="left", padx=24, pady=12)
        tk.Label(hdr_inner, text=r.get("name", "Unknown"),
                 font=(FONT_FAMILY, 15, "bold"),
                 bg=PRIMARY, fg="white").pack(anchor="w")
        tk.Label(hdr_inner, text=r.get("_file", ""),
                 font=FONT_SMALL, bg=PRIMARY, fg=TEXT_SUBTLE).pack(anchor="w")

        score    = r.get("score", 0)
        decision = r.get("decision", "Reject")
        is_sel   = str(decision).strip().lower() == "select"

        badge = tk.Frame(hdr, bg=PRIMARY)
        badge.pack(side="right", padx=24, pady=14)
        tk.Label(
            badge,
            text=("✓ " if is_sel else "✕ ") + decision,
            font=FONT_BOLD,
            bg=SUCCESS_BG if is_sel else DANGER_BG,
            fg=SUCCESS if is_sel else DANGER,
            padx=14, pady=6,
        ).pack(side="right")

        # Body
        body = tk.Frame(win, bg=BG)
        body.pack(fill="both", expand=True, padx=20, pady=16)

        # Stat row
        stats = tk.Frame(body, bg=BG)
        stats.pack(fill="x", pady=(0, 14))

        def stat(parent, label, value, fg=TEXT):
            box = tk.Frame(parent, bg=CARD,
                           highlightbackground=BORDER, highlightthickness=1)
            box.pack(side="left", fill="x", expand=True, padx=(0, 10))
            tk.Label(box, text=label, font=FONT_LABEL,
                     bg=CARD, fg=TEXT_MUTED).pack(anchor="w", padx=14, pady=(10, 0))
            tk.Label(box, text=value, font=(FONT_FAMILY, 14, "bold"),
                     bg=CARD, fg=fg).pack(anchor="w", padx=14, pady=(0, 10))

        stat(stats, "TOTAL EXPERIENCE",    fmt_exp(r.get("total_experience")))
        stat(stats, "RELEVANT EXPERIENCE", fmt_exp(r.get("relevant_experience")))
        stat(stats, "MATCH SCORE", f"{score} / 10",
             fg=SUCCESS if is_sel else DANGER)

        def section(title: str, items, color):
            wrap = tk.Frame(body, bg=CARD,
                            highlightbackground=BORDER, highlightthickness=1)
            wrap.pack(fill="x", pady=(0, 10))
            head = tk.Frame(wrap, bg=CARD)
            head.pack(fill="x", padx=14, pady=(10, 4))
            tk.Label(head, text=title, font=FONT_BOLD,
                     bg=CARD, fg=color).pack(anchor="w")
            for it in (items or ["—"]):
                row = tk.Frame(wrap, bg=CARD)
                row.pack(fill="x", padx=14, pady=2)
                tk.Label(row, text="•", font=FONT_BOLD,
                         bg=CARD, fg=color).pack(side="left", padx=(0, 6))
                tk.Label(row, text=it, font=FONT_BODY, bg=CARD, fg=TEXT,
                         wraplength=640, justify="left",
                         anchor="w").pack(side="left", fill="x", expand=True)
            tk.Frame(wrap, bg=CARD, height=8).pack()

        # Profile card: education + organisations
        edu_lines = []
        seen = set()
        for e in (r.get("education") or []):
            if not isinstance(e, dict):
                continue
            deg  = str(e.get("degree", "") or "").strip()
            coll = str(e.get("college", "") or "").strip()
            yr   = e.get("year")
            yr_s = str(int(yr)) if isinstance(yr, (int, float)) and yr else (
                str(yr).strip() if yr else ""
            )
            parts = [p for p in (deg, coll, yr_s) if p]
            if not parts:
                continue
            line = ", ".join(parts)
            k = line.lower()
            if k in seen:
                continue
            seen.add(k)
            edu_lines.append(line)

        current_org = str(r.get("current_organization") or "").strip()
        past_orgs_raw = r.get("past_organizations") or []
        past_orgs = []
        seen_o = set()
        excl = current_org.lower()
        for o in past_orgs_raw:
            s = str(o or "").strip()
            if not s:
                continue
            k = s.lower()
            if k in seen_o or k == excl:
                continue
            seen_o.add(k)
            past_orgs.append(s)

        prof = tk.Frame(body, bg=CARD,
                        highlightbackground=BORDER, highlightthickness=1)
        prof.pack(fill="x", pady=(0, 10))
        tk.Label(prof, text="Profile", font=FONT_BOLD,
                 bg=CARD, fg=TEXT).pack(anchor="w", padx=14, pady=(10, 4))

        def info_row(label: str, value: str, mono: bool = False):
            row = tk.Frame(prof, bg=CARD)
            row.pack(fill="x", padx=14, pady=3)
            tk.Label(row, text=label, font=FONT_LABEL,
                     bg=CARD, fg=TEXT_MUTED, width=12,
                     anchor="w").pack(side="left")
            tk.Label(row, text=value or "—",
                     font=FONT_BODY, bg=CARD, fg=TEXT,
                     wraplength=560, justify="left",
                     anchor="w").pack(side="left", fill="x", expand=True)

        info_row("EDUCATION",
                 "\n".join(edu_lines) if edu_lines else "—")
        info_row("CURRENT ORG", current_org)
        info_row("PAST ORGS",
                 ", ".join(past_orgs) if past_orgs else "—")
        tk.Frame(prof, bg=CARD, height=8).pack()

        section("Strengths",  r.get("highlights"), SUCCESS)
        section("Concerns",   r.get("lowlights"),  DANGER)

        footer = tk.Frame(win, bg=BG)
        footer.pack(fill="x", padx=20, pady=(0, 16))
        HoverButton(footer, "Close", win.destroy,
                    bg=NEUTRAL_BTN, fg=TEXT, hover=NEUTRAL_HOVER,
                    font=FONT_BOLD, padx=20, pady=8).pack(side="right")

    # ── Export ────────────────────────────────────────────────────────────────

    def _export_csv(self):
        if not self.results:
            messagebox.showinfo("No Data", "Run the analysis first.")
            return
        path = filedialog.asksaveasfilename(
            defaultextension=".csv",
            filetypes=[("CSV Files", "*.csv")],
            initialfile="resume_analysis",
        )
        if not path:
            return
        import csv
        with open(path, "w", newline="", encoding="utf-8-sig") as f:
            w = csv.writer(f)
            w.writerow([
                "#", "Name", "Education", "Current Organization",
                "Past Organizations", "Total Experience", "Relevant Experience",
                "Highlights", "Lowlights", "Score", "Decision",
            ])
            for i, r in enumerate(self.results, 1):
                current_org = str(r.get("current_organization") or "").strip()
                w.writerow([
                    i, r.get("name", ""),
                    fmt_education(r.get("education")),
                    current_org,
                    fmt_orgs(r.get("past_organizations"), exclude=current_org),
                    fmt_exp(r.get("total_experience")),
                    fmt_exp(r.get("relevant_experience")),
                    " | ".join(r.get("highlights", [])),
                    " | ".join(r.get("lowlights", [])),
                    r.get("score", ""),
                    r.get("decision", ""),
                ])
        messagebox.showinfo("Exported", f"CSV saved to:\n{path}")

    def _export_excel(self):
        if not self.results:
            messagebox.showinfo("No Data", "Run the analysis first.")
            return
        try:
            import openpyxl
            from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
        except ImportError:
            messagebox.showerror("Error",
                                 "openpyxl not available. Please use CSV export.")
            return

        path = filedialog.asksaveasfilename(
            defaultextension=".xlsx",
            filetypes=[("Excel Files", "*.xlsx")],
            initialfile="resume_analysis",
        )
        if not path:
            return

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Resume Analysis"

        h_fill   = PatternFill("solid", fgColor="0F172A")
        h_font   = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
        thin     = Side(style="thin", color="E2E8F0")
        border   = Border(left=thin, right=thin, top=thin, bottom=thin)
        c_align  = Alignment(horizontal="center", vertical="top", wrap_text=True)
        l_align  = Alignment(horizontal="left",   vertical="top", wrap_text=True)

        headers  = ["#", "Name", "Education", "Current Organization",
                    "Past Organizations", "Total Experience", "Relevant Experience",
                    "Highlights", "Lowlights", "Score", "Decision"]
        for col, h in enumerate(headers, 1):
            cell           = ws.cell(row=1, column=col, value=h)
            cell.fill      = h_fill
            cell.font      = h_font
            cell.alignment = c_align
            cell.border    = border

        sel_fill = PatternFill("solid", fgColor="D1FAE5")
        rej_fill = PatternFill("solid", fgColor="FEE2E2")
        alt_fill = PatternFill("solid", fgColor="F8FAFC")

        # Columns to center-align (1-indexed): #, Total Exp, Rel Exp, Score, Decision
        center_cols = {1, 6, 7, 10, 11}

        for ri, r in enumerate(self.results, 2):
            decision = str(r.get("decision", "Reject")).strip().lower()
            if decision == "select":
                row_fill = sel_fill
            elif ri % 2 == 0:
                row_fill = alt_fill
            else:
                row_fill = rej_fill

            current_org = str(r.get("current_organization") or "").strip()
            # Education: each entry on its own line for readability
            edu_lines, seen = [], set()
            for e in (r.get("education") or []):
                if not isinstance(e, dict):
                    continue
                deg  = str(e.get("degree", "") or "").strip()
                coll = str(e.get("college", "") or "").strip()
                yr   = e.get("year")
                yr_s = str(int(yr)) if isinstance(yr, (int, float)) and yr else (
                    str(yr).strip() if yr else ""
                )
                parts = [p for p in (deg, coll, yr_s) if p]
                if not parts:
                    continue
                line = ", ".join(parts)
                k = line.lower()
                if k not in seen:
                    seen.add(k)
                    edu_lines.append(line)

            past_lines, seen_o = [], set()
            excl = current_org.lower()
            for o in (r.get("past_organizations") or []):
                s = str(o or "").strip()
                if not s:
                    continue
                k = s.lower()
                if k in seen_o or k == excl:
                    continue
                seen_o.add(k)
                past_lines.append(s)

            values = [
                ri - 1,
                r.get("name", ""),
                "\n".join(edu_lines) if edu_lines else "",
                current_org,
                "\n".join(past_lines) if past_lines else "",
                fmt_exp(r.get("total_experience")),
                fmt_exp(r.get("relevant_experience")),
                "\n".join(f"• {h}" for h in r.get("highlights", [])),
                "\n".join(f"• {l}" for l in r.get("lowlights",  [])),
                r.get("score", ""),
                r.get("decision", ""),
            ]
            for ci, val in enumerate(values, 1):
                cell           = ws.cell(row=ri, column=ci, value=val)
                cell.fill      = row_fill
                cell.border    = border
                cell.alignment = c_align if ci in center_cols else l_align

        for col, width in zip(
            "ABCDEFGHIJK",
            [5, 24, 36, 24, 32, 16, 16, 42, 32, 9, 13],
        ):
            ws.column_dimensions[col].width = width

        ws.freeze_panes = "A2"

        wb.save(path)
        messagebox.showinfo("Exported", f"Excel file saved to:\n{path}")


# ─── Entry point ──────────────────────────────────────────────────────────────

def main():
    app = ResumeAnalyzerApp()
    app.mainloop()


if __name__ == "__main__":
    main()
